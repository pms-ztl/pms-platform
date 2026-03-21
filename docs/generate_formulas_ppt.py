"""
Generate PMS Backend Formulas Presentation (.pptx)
Each formula gets its own slide with: What it does, Parameters, Formula, Result/Threshold.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──
NAVY    = RGBColor(0x0F, 0x17, 0x2A)
VIOLET  = RGBColor(0x6D, 0x28, 0xD9)
BLUE    = RGBColor(0x1E, 0x3A, 0x5F)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xE2, 0xE8, 0xF0)
GRAY    = RGBColor(0x94, 0xA3, 0xB8)
TEAL    = RGBColor(0x14, 0xB8, 0xA6)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
RED     = RGBColor(0xEF, 0x44, 0x44)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
SLATE   = RGBColor(0x33, 0x40, 0x55)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = 'Calibri'
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_size=13, color=WHITE, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold, txt_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = txt_color if txt_color else color
        run.font.bold = is_bold
        run.font.name = 'Calibri'
    return txBox


# ═══════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, NAVY)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
             'PMS Platform', 44, VIOLET, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(0.8),
             'Backend Mathematical Engine', 32, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.6),
             '17 Core Formulas  +  6 Composite Systems  +  CPIS 8-Dimension Scoring', 16, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             'How Every Number in the System is Calculated', 18, TEAL, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.4),
             'Confidential  |  March 2026', 12, GRAY, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 2: OVERVIEW
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             'Formula Architecture Overview', 28, VIOLET, True)

# Left column - Core Formulas
add_rect(slide, Inches(0.6), Inches(1.2), Inches(3.8), Inches(5.8), SLATE)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(3.4), Inches(0.5),
             'LAYER 1: Core Statistical (17)', 16, TEAL, True)
core_lines = [
    'Clamp Percentage', 'Arithmetic Mean', 'Weighted Mean',
    'Weighted Harmonic Mean', 'Variance & Std Deviation',
    'Sample Std Dev (Bessel)', 'Z-Score Normalization',
    'Sigmoid Function', 'Bounded Sigmoid',
    'Pearson Correlation', 'Linear Regression',
    'EWMA (Exp. Smoothing)', 'Shannon Entropy',
    'Percentile Rank', 'Bayesian Estimation', 'Gini Coefficient',
]
for i, name in enumerate(core_lines):
    add_text_box(slide, Inches(1.0), Inches(1.85 + i * 0.24), Inches(3.2), Inches(0.3),
                 f'{i+1}. {name}', 10.5, LIGHT)

# Middle column - Composite
add_rect(slide, Inches(4.7), Inches(1.2), Inches(3.8), Inches(5.8), SLATE)
add_text_box(slide, Inches(4.9), Inches(1.3), Inches(3.4), Inches(0.5),
             'LAYER 2: Composite Systems (6)', 16, AMBER, True)
composites = [
    ('Goal Score', 'Completion + Quality + Timeliness + Efficiency'),
    ('Performance Score', 'Goals + Reviews + Feedback + Attendance + Collab'),
    ('Team Analytics', 'Mean + StdDev + Entropy + Gini + Velocity'),
    ('Goal Risk', 'Schedule + Velocity + Dependency + Complexity'),
    ('Review Calibration', 'Z-Score normalization per reviewer'),
    ('Disparate Impact', '4/5ths rule fairness analysis'),
]
for i, (name, desc) in enumerate(composites):
    y = Inches(1.9 + i * 0.7)
    add_text_box(slide, Inches(5.0), y, Inches(3.4), Inches(0.3),
                 name, 13, WHITE, True)
    add_text_box(slide, Inches(5.0), y + Inches(0.25), Inches(3.4), Inches(0.3),
                 desc, 10, GRAY)

# Right column - CPIS
add_rect(slide, Inches(8.8), Inches(1.2), Inches(4), Inches(5.8), SLATE)
add_text_box(slide, Inches(9.0), Inches(1.3), Inches(3.6), Inches(0.5),
             'LAYER 3: CPIS (8 Dimensions)', 16, GREEN, True)
cpis_dims = [
    ('GAI 25%', 'Goal Attainment Index'),
    ('RQS 20%', 'Review Quality Score'),
    ('FSI 12%', 'Feedback Sentiment Index'),
    ('CIS 10%', 'Collaboration Impact Score'),
    ('CRI 10%', 'Consistency & Reliability'),
    ('GTS 8%', 'Growth Trajectory Score'),
    ('EQS 8%', 'Evidence Quality Score'),
    ('III 7%', 'Initiative & Innovation'),
]
for i, (code, name) in enumerate(cpis_dims):
    y = Inches(1.9 + i * 0.6)
    add_text_box(slide, Inches(9.1), y, Inches(1.6), Inches(0.3),
                 code, 13, TEAL, True)
    add_text_box(slide, Inches(10.6), y, Inches(2.0), Inches(0.3),
                 name, 11, LIGHT)


# ═══════════════════════════════════════════════════════════
# FORMULA SLIDE TEMPLATE
# ═══════════════════════════════════════════════════════════
def formula_slide(number, title, what_it_does, parameters, formula_text, result_thresholds, where_used):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Header bar
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), RGBColor(0x14, 0x1E, 0x30))
    add_text_box(slide, Inches(0.5), Inches(0.05), Inches(1), Inches(0.5),
                 f'#{number}', 28, VIOLET, True)
    add_text_box(slide, Inches(1.5), Inches(0.1), Inches(10), Inches(0.5),
                 title, 26, WHITE, True)
    add_text_box(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.4),
                 where_used, 12, GRAY, False)

    # LEFT: What + Parameters
    # What it does box
    add_rect(slide, Inches(0.4), Inches(1.3), Inches(5.8), Inches(1.5), SLATE, TEAL)
    add_text_box(slide, Inches(0.6), Inches(1.35), Inches(5.4), Inches(0.35),
                 'WHAT IT MEASURES', 11, TEAL, True)
    add_text_box(slide, Inches(0.6), Inches(1.7), Inches(5.4), Inches(1.0),
                 what_it_does, 13, LIGHT)

    # Parameters box
    param_height = max(2.0, len(parameters) * 0.35 + 0.5)
    add_rect(slide, Inches(0.4), Inches(3.0), Inches(5.8), Inches(param_height), SLATE, AMBER)
    add_text_box(slide, Inches(0.6), Inches(3.05), Inches(5.4), Inches(0.35),
                 'PARAMETERS', 11, AMBER, True)
    for i, param in enumerate(parameters):
        add_text_box(slide, Inches(0.6), Inches(3.45 + i * 0.32), Inches(5.4), Inches(0.3),
                     f'  {param}', 11.5, LIGHT)

    # RIGHT: Formula + Results
    # Formula box
    add_rect(slide, Inches(6.5), Inches(1.3), Inches(6.4), Inches(2.0), RGBColor(0x1A, 0x10, 0x30), VIOLET)
    add_text_box(slide, Inches(6.7), Inches(1.35), Inches(6.0), Inches(0.35),
                 'FORMULA', 11, VIOLET, True)
    add_text_box(slide, Inches(6.7), Inches(1.75), Inches(6.0), Inches(1.4),
                 formula_text, 14, WHITE, True)

    # Result / Threshold box
    result_height = max(2.0, len(result_thresholds) * 0.32 + 0.6)
    add_rect(slide, Inches(6.5), Inches(3.5), Inches(6.4), Inches(result_height), SLATE, GREEN)
    add_text_box(slide, Inches(6.7), Inches(3.55), Inches(6.0), Inches(0.35),
                 'RESULT & THRESHOLDS', 11, GREEN, True)
    for i, line in enumerate(result_thresholds):
        add_text_box(slide, Inches(6.7), Inches(3.95 + i * 0.3), Inches(6.0), Inches(0.3),
                     f'  {line}', 11.5, LIGHT)

    return slide


# ═══════════════════════════════════════════════════════════
# ALL 17 CORE FORMULAS
# ═══════════════════════════════════════════════════════════

formula_slide(
    1, 'Clamp Percentage',
    'Bounds any number to a valid percentage range (0-100). Acts as a safety net ensuring no score goes below 0 or above 100.',
    ['x = any input number', 'min = 0 (lower bound)', 'max = 100 (upper bound)'],
    'clampPct(x) = max(0, min(100, x))',
    ['Output: Always 0 to 100', 'Example: clampPct(115) = 100', 'Example: clampPct(-5) = 0', 'Example: clampPct(78) = 78'],
    'Used in: Every score calculation as final safety bound'
)

formula_slide(
    2, 'Arithmetic Mean (Simple Average)',
    'Calculates the basic average of a set of numbers. Sum all values, divide by count. Treats every value equally.',
    ['x1, x2, ..., xn = set of values', 'n = count of values'],
    'Mean = (x1 + x2 + ... + xn) / n',
    ['Output: Single number (same unit as inputs)', 'Example: mean(80, 90, 70) = 80', 'No threshold — used as building block'],
    'Used in: Team average scores, department averages'
)

formula_slide(
    3, 'Weighted Arithmetic Mean',
    'Average where some values count more than others. Like how a final exam counts more than homework. Each value is multiplied by its importance weight.',
    ['xi = each value', 'wi = weight (importance) of each value', 'Sum of weights does NOT need to equal 1'],
    'Weighted Mean = Sum(wi * xi) / Sum(wi)',
    ['Output: Single weighted average', 'Example: Goals(80)x0.4 + Reviews(90)x0.3 + Feedback(70)x0.3', '= (32 + 27 + 21) / 1.0 = 80', 'Higher weights = more influence on final score'],
    'Used in: CPIS dimension aggregation, goal composite scores'
)

formula_slide(
    4, 'Weighted Harmonic Mean',
    'A special average that is better for rates and ratios. Unlike arithmetic mean, it prevents one very high score from masking many low ones. More "pessimistic" — pulls toward lower values.',
    ['xi = each rating value (must be > 0)', 'wi = weight per rating (e.g., reviewer trust)', 'Cannot handle zero values'],
    'WHM = Sum(wi) / Sum(wi / xi)',
    ['Output: Single number, always <= arithmetic mean', 'Example: ratings [5, 2, 2] with equal weights', '  Arithmetic mean = 3.0', '  Harmonic mean = 2.5 (closer to reality)', 'Prevents inflated averages'],
    'Used in: Review Quality Score (RQS) — combining trust-weighted ratings'
)

formula_slide(
    5, 'Variance & Standard Deviation',
    'Measures how spread out scores are from the average. Low = everyone is similar. High = big differences. Standard deviation is the square root of variance (same unit as data).',
    ['xi = each data point', 'mu = mean of all data points', 'n = count of data points'],
    'Variance = Sum((xi - mu)^2) / n\nStd Dev = sqrt(Variance)',
    ['Output: Number >= 0', 'StdDev = 0 means all scores identical', 'StdDev = 5 on a 0-100 scale = very consistent', 'StdDev = 20 = wide spread, big gaps in team'],
    'Used in: Team analytics, rating consistency, calibration'
)

formula_slide(
    6, 'Sample Standard Deviation (Bessel\'s Correction)',
    'When working with a sample (not full data), we divide by (n-1) instead of n. This avoids underestimating the true spread. Important when team size is small.',
    ['xi = each data point', 'x_bar = sample mean', 'n = sample size (must be >= 2)'],
    'Sample StdDev = sqrt(Sum((xi - x_bar)^2) / (n - 1))',
    ['Output: Slightly larger than regular StdDev', 'Critical when n < 30 (small teams)', 'Falls back to population StdDev when n = 1', 'Ensures fair spread estimation for small groups'],
    'Used in: Team comparisons when team sizes differ'
)

formula_slide(
    7, 'Z-Score Normalization',
    'Measures how many standard deviations a value is above or below average. Like grading on a curve — a Z-score of +2 means "way above average" regardless of the original scale.',
    ['x = individual value', 'mu = group mean (average)', 'sigma = group standard deviation'],
    'Z = (x - mu) / sigma',
    ['Output: Usually -3 to +3', 'Z = 0: Exactly average', 'Z = +1: One std dev above average (top ~16%)', 'Z = +2: Two std dev above (top ~2%)', 'Z < -1: Below average (flag for support)'],
    'Used in: Peer comparison, review calibration, anomaly detection'
)

formula_slide(
    8, 'Sigmoid Function',
    'Maps ANY number to a smooth S-curve between 0 and 1. Small inputs give near-0, large inputs give near-1, with a rapid transition zone in the middle. Perfect for converting raw counts to scores.',
    ['x = input value', 'k = steepness (how sharp the transition)', 'x0 = midpoint (where output = 0.5)'],
    'sigmoid(x) = 1 / (1 + e^(-k * (x - x0)))',
    ['Output: 0 to 1 (smooth curve)', 'At x = x0: output = 0.5 (midpoint)', 'k = 0.2: gentle curve', 'k = 3: sharp threshold-like curve', 'Multiply by 100 for percentage'],
    'Used in: Collaboration scores, risk scores, initiative metrics'
)

formula_slide(
    9, 'Bounded Sigmoid',
    'Sigmoid that maps to a custom range instead of 0-1. For example, map collaboration activity count to 0-100 scale, or timeliness to 0.6-1.4 multiplier.',
    ['x = input value', 'min = minimum output', 'max = maximum output', 'k = steepness', 'x0 = midpoint'],
    'bounded(x) = min + (max - min) * sigmoid(x, k, x0)',
    ['Output: min to max (configurable)', 'Example: 3 cross-team goals, k=0.5, x0=3', '  -> sigmoid = 0.5 -> score = 50/100', 'Example: 8 goals -> sigmoid = 0.92 -> score = 92'],
    'Used in: CPIS collaboration channels, timeliness factor, initiative scoring'
)

formula_slide(
    10, 'Pearson Correlation',
    'Measures the strength of a straight-line relationship between two variables. Do they move together (+1), move opposite (-1), or not related (0)?',
    ['X = first variable (e.g., attendance rates)', 'Y = second variable (e.g., performance scores)', 'n = number of paired observations'],
    'r = Sum((xi-x_bar)(yi-y_bar)) / sqrt(Sum((xi-x_bar)^2) * Sum((yi-y_bar)^2))',
    ['Output: -1 to +1', 'r = +0.8 to +1.0: Strong positive (they rise together)', 'r = +0.3 to +0.8: Moderate positive', 'r = -0.3 to +0.3: Weak / no relationship', 'r = -1.0 to -0.3: Negative (one rises, other falls)'],
    'Used in: Analytics insights, finding hidden performance drivers'
)

formula_slide(
    11, 'Linear Regression (Trend Line)',
    'Fits a straight line through data points to show the trend. Positive slope = improving over time. Also gives R-squared: how well the line fits the data (0 = random, 1 = perfect fit).',
    ['x = time points (e.g., quarter 1, 2, 3)', 'y = scores at each time point', 'n = number of data points (need >= 2)'],
    'y = mx + b\nwhere m = slope, b = intercept\nR^2 = goodness of fit (0 to 1)',
    ['m > 0: Performance improving', 'm = 0: Stable/flat', 'm < 0: Performance declining', 'R^2 > 0.7: Strong trend (reliable prediction)', 'R^2 < 0.3: Weak trend (noisy data)'],
    'Used in: Growth Trajectory Score, team velocity trends, score predictions'
)

formula_slide(
    12, 'EWMA (Exponentially Weighted Moving Average)',
    'A "smart average" that gives more weight to recent data. Like memory that gradually fades — yesterday matters more than last year. Recent events shape the score more.',
    ['xt = current value', 'EWMA(t-1) = previous EWMA', 'alpha = smoothing factor (0-1)', '  alpha=0.3 = moderate recency bias', '  alpha=0.4 = stronger recency bias'],
    'EWMA(t) = alpha * x(t) + (1 - alpha) * EWMA(t-1)',
    ['Output: Smoothed trend value', 'alpha = 0.3: 30% current, 70% history', 'alpha = 0.4: 40% current, 60% history', 'Higher alpha = more reactive to recent changes', 'Lower alpha = more stable, less reactive'],
    'Used in: Feedback Sentiment (alpha=0.35), Quality Scores (alpha=0.4)'
)

formula_slide(
    13, 'Shannon Entropy',
    'Measures how evenly distributed ratings are. If everyone gets the same rating = 0 (no diversity). If ratings are perfectly spread across all levels = 1 (maximum diversity). Detects "rating clustering."',
    ['pi = proportion in each rating bucket', 'k = number of possible ratings (e.g., 5)', 'Normalized to 0-1 scale'],
    'H = -Sum(pi * log2(pi)) / log2(k)',
    ['Output: 0 to 1', 'H = 0: All same rating (everyone got a 4)', 'H = 0.5: Some clustering', 'H = 1.0: Perfect spread across all ratings', 'H < 0.3: WARNING — manager may be rating lazily'],
    'Used in: Team analytics, calibration health checks'
)

formula_slide(
    14, 'Percentile Rank',
    'Shows what percentage of people scored lower. The 82nd percentile means "better than 82% of peers." Uses interpolation for ties to avoid unfair clumping.',
    ['score = individual\'s score', 'all_scores = sorted list of all scores', 'n = total count'],
    'Percentile = (count_below + 0.5 * count_equal) / n * 100',
    ['Output: 0 to 100', '90th percentile: Top 10% performer', '75th percentile: Above average', '50th percentile: Exactly average', '25th percentile: Below average', '10th percentile: Bottom performer'],
    'Used in: Leaderboards, performance ranking, derived rating (1-5 stars)'
)

formula_slide(
    15, 'Bayesian Estimation',
    'When individual data is scarce (new employee), it "borrows strength" from the group average. As more data comes in, the estimate gradually shifts to trust individual data more. Prevents unfair scoring of new hires.',
    ['observed_mean = person\'s actual average', 'n = person\'s data count', 'prior_mean = department/org average', 'prior_weight = how much to trust group (shrinks as n grows)'],
    'Bayesian = (prior_mean * prior_weight + observed_mean * n) / (prior_weight + n)',
    ['Output: Blended score between individual and group', 'New hire (n=1): Heavily influenced by dept avg', 'Experienced (n=20): Almost entirely their own data', 'prior_weight = max(0, 5 - n * 0.5)', 'At n >= 10: prior_weight = 0 (fully individual)'],
    'Used in: CPIS fairness adjustment, new employee scoring'
)

formula_slide(
    16, 'Gini Coefficient',
    'Measures inequality in distribution — like wealth inequality. Are performance rewards concentrated in a few top performers, or distributed fairly? 0 = perfect equality, 1 = maximum inequality.',
    ['xi = sorted scores (ascending)', 'n = count of scores', 'Requires sorted input'],
    'G = (2 * Sum(i * xi)) / (n * Sum(xi)) - (n + 1) / n',
    ['Output: 0 to 1', 'G < 0.2: Very equal distribution (healthy team)', 'G = 0.2-0.4: Moderate inequality (normal)', 'G = 0.4-0.6: High inequality (concerning)', 'G > 0.6: Extreme inequality (investigate why)'],
    'Used in: Team equity analysis, compensation fairness'
)

formula_slide(
    17, 'Disparate Impact (4/5ths Rule)',
    'A legal fairness test used in employment law. If one group\'s score falls below 80% of the top group\'s score, it flags potential discrimination. This is the same test used by courts.',
    ['raw_score = individual\'s CPIS score', 'dept_avg = department average score', 'threshold = 0.8 (the 4/5ths rule)'],
    'Disparate Impact Ratio = raw_score / dept_avg\nFLAG if ratio < 0.8',
    ['Output: Ratio (typically 0.5 to 1.5)', 'Ratio >= 0.8: PASS — no disparate impact', 'Ratio < 0.8: FLAG — potential bias detected', 'Triggers manual review by HR', 'Legally defensible metric (EEOC standard)'],
    'Used in: CPIS fairness analysis, bias detection dashboard'
)


# ═══════════════════════════════════════════════════════════
# COMPOSITE SYSTEMS (6 slides)
# ═══════════════════════════════════════════════════════════

# Divider slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
             'COMPOSITE SYSTEMS', 40, AMBER, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
             '6 Multi-Formula Calculations Built on the Core 17', 18, GRAY, False, PP_ALIGN.CENTER)

formula_slide(
    18, 'Goal Completion Score (Composite)',
    'Combines 4 factors into one goal score: How much is done (completion), how well (quality), how fast (timeliness), and effort vs difficulty (efficiency). Each factor uses different core formulas.',
    ['Completion = weighted avg of task progress', 'Quality = EWMA of quality ratings / 5 * 100', 'Timeliness = sigmoid(-avgDaysLate) mapped to 0-100', 'Efficiency = clamp(quality/complexity, 0.5-1.5)',
     'Weights: Completion 50%, Quality 30%, Timeliness 20%'],
    'GoalScore = (0.50 * Completion + 0.30 * Quality + 0.20 * Timeliness) * Efficiency',
    ['Output: 0 to ~150 (clamped to 100)', '90-100: Exceptional goal delivery', '70-89: Strong performance', '50-69: Meets expectations', 'Below 50: Needs improvement', 'Also calculates: velocity, predicted completion, risk'],
    'Used in: CPIS Goal Attainment Index (GAI), goal dashboards'
)

formula_slide(
    19, 'Individual Performance Score (Composite)',
    'The overall performance score combining 5 data sources. Uses dynamic weights — if a component has no data, its weight redistributes to others. No one is penalized for missing data.',
    ['Goals = weighted mean of goal composite scores', 'Reviews = harmonic mean (trust-weighted, bias-corrected)', 'Feedback = EWMA of sentiment mapped 0-100',
     'Attendance = attendance rate * 100', 'Collaboration = sigmoid-based channel scores',
     'Default weights: Goals 40%, Reviews 30%, Feedback 10%, Attend 10%, Collab 10%'],
    'PerfScore = Sum(dynamic_weight_i * component_i)\nDerived Rating: Percentile-based 1-5 stars',
    ['Output: 0-100 score + 1-5 star rating', 'Top 10% -> 5 stars', 'Next 20% -> 4 stars', 'Middle 40% -> 3 stars', 'Next 20% -> 2 stars', 'Bottom 10% -> 1 star',
     'Also outputs: confidence score (0-1)'],
    'Used in: Employee profiles, rankings, promotion decisions'
)

formula_slide(
    20, 'Team Analytics (Composite)',
    'Provides a statistical profile of a team: average performance, spread, fairness of distribution, and whether the team is improving or declining. Uses 5 different statistical measures.',
    ['scores = array of all team member scores', 'historical = past team averages over time', 'Minimum 2 members required'],
    'TeamProfile = {\n  avg: mean(scores),\n  spread: stddev(scores),\n  entropy: shannon(rating_distribution),\n  gini: gini_coefficient(scores),\n  trend: linear_regression(historical).slope\n}',
    ['Average: team health indicator', 'Spread < 10: Consistent team', 'Spread > 25: High variance (investigate)', 'Entropy < 0.3: Rating clustering (lazy manager?)',
     'Gini > 0.4: Inequality concern', 'Trend slope > 0: Team improving'],
    'Used in: Team dashboards, manager analytics, HR oversight'
)

formula_slide(
    21, 'Goal Risk Assessment (Composite)',
    'Predicts the probability of a goal failing. Combines 4 risk factors into one 0-100 risk score. Higher = more likely to miss. Triggers early warnings for at-risk goals.',
    ['Schedule = sigmoid of (expected - actual progress)', 'Velocity = sigmoid of (required pace / current pace)',
     'Dependency = mean + max of dependent goal risks', 'Complexity = (complexity/5) * remaining%',
     'Weights: Schedule 40%, Velocity 30%, Dependency 15%, Complexity 15%'],
    'Risk = 0.40*Schedule + 0.30*Velocity + 0.15*Dependency + 0.15*Complexity',
    ['Output: 0-100 risk score', 'CRITICAL >= 75: Likely to fail, escalate now', 'HIGH 50-74: At risk, needs intervention', 'MEDIUM 25-49: Monitor closely',
     'LOW < 25: On track'],
    'Used in: Goal dashboards, manager alerts, risk heatmaps'
)

formula_slide(
    22, 'Review Rating Calibration',
    'Removes reviewer bias using Z-score normalization. Some managers rate everyone high (lenient), others rate everyone low (strict). This formula "translates" ratings to a common scale.',
    ['rating = original rating given (1-5)', 'reviewer_mean = this reviewer\'s average', 'reviewer_stddev = this reviewer\'s spread',
     'global_mean = everyone\'s average', 'global_stddev = everyone\'s spread'],
    'z = (rating - reviewer_mean) / reviewer_stddev\ncalibrated = z * global_stddev + global_mean\nResult = clamp(calibrated, 1, 5)',
    ['Output: Calibrated rating (1-5)', 'Strict reviewer\'s 3.0 -> might become 3.8', 'Lenient reviewer\'s 4.5 -> might become 3.7', 'Fair reviewers: minimal change', 'Ensures cross-team fairness'],
    'Used in: CPIS Review Quality Score, calibration sessions'
)


# ═══════════════════════════════════════════════════════════
# CPIS DIMENSION SLIDES (8 slides)
# ═══════════════════════════════════════════════════════════

# Divider
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1),
             'CPIS: 8 Dimensions', 40, GREEN, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.6),
             'Comprehensive Performance Intelligence Score', 20, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
             'Final Score = Sum(Dimension * Weight) * TenureFactor * FairnessAdjustment', 16, TEAL, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.6),
             'Output: 0-100 Score | Letter Grade (A+ to F) | Star Rating (1-5) | Rank Label | Confidence Interval', 14, GRAY, False, PP_ALIGN.CENTER)

formula_slide(
    'D1', 'CPIS Dimension 1: Goal Attainment Index (GAI) — 25%',
    'Measures how well goals are being completed. Factors in completion progress, priority level, timeliness, and alignment to company objectives. The heaviest weighted dimension.',
    ['Gi = goal progress (0-100)', 'Wi = goal weight (importance)', 'Pi = priority multiplier: LOW=0.85, MED=1.0, HIGH=1.15, CRITICAL=1.35',
     'Ti = timeliness: boundedSigmoid(-daysLate, 0.6, 1.4)', 'Ai = alignment bonus: 1 + depth*0.03 (max 1.15)'],
    'GAI = Sum(Gi * Wi * Pi * Ti * Ai) / Sum(Wi)',
    ['Output: 0-100', '90+: Exceptional goal delivery', '70-89: Strong attainment', '50-69: Meets expectations', '< 50: Below expectations',
     'Sub-metrics: completion rate, on-time rate, aligned goals count'],
    'Weight: 25% of final CPIS score'
)

formula_slide(
    'D2', 'CPIS Dimension 2: Review Quality Score (RQS) — 20%',
    'Combines all review ratings with bias correction. Uses weighted harmonic mean (fairer than arithmetic mean). Adjusts for reviewer trust, review type, and detected bias.',
    ['Ri = calibrated rating (1-5)', 'Bi = bias score (0-1), penalizes up to 50%', 'Ti = reviewer trust (0-1)',
     'TypeW: Manager=1.5, 360=1.3, Peer=1.0, Upward=0.9, Self=0.5'],
    'RQS = WHM(Ri * (1 - Bi), Ti * TypeWi) * 100 / 5\nDefault: 50 if no reviews (neutral)',
    ['Output: 0-100', 'Manager reviews count 3x more than self-reviews', 'Bias > 0.3 triggers penalty', 'Default 50 = not penalized for no reviews',
     'Sub-metrics: raw avg, calibrated avg, trust-weighted avg'],
    'Weight: 20% of final CPIS score'
)

formula_slide(
    'D3', 'CPIS Dimension 3: Feedback Sentiment Index (FSI) — 12%',
    'Tracks the tone of feedback received over time. Uses EWMA so recent feedback counts more. Quality bonus for well-tagged feedback (with skill/value tags).',
    ['Si = sentiment score (0-1, from NLP analysis)', 'Qi = quality multiplier: base 1.0, +0.1 for skill tags, +0.1 for value tags',
     'alpha = 0.35 (recency weight)'],
    'FSI = EWMA(Si * Qi, alpha=0.35) * 100\nDefault: 50 if no feedback',
    ['Output: 0-100', '80+: Overwhelmingly positive feedback', '60-79: Mostly positive', '40-59: Mixed or neutral',
     '< 40: Mostly negative feedback', 'Sub-metrics: praise ratio, feedback count, quality multiplier avg'],
    'Weight: 12% of final CPIS score'
)

formula_slide(
    'D4', 'CPIS Dimension 4: Collaboration Impact Score (CIS) — 10%',
    'Measures how much someone collaborates across the organization. Tracks 6 channels of collaboration, each converted to 0-100 via sigmoid curves.',
    ['6 channels, each via boundedSigmoid(count, 0, 100, k, x0):', '  Cross-functional goals (w=0.20, x0=3)',
     '  Feedback given (w=0.15, x0=5)', '  Feedback received (w=0.15, x0=5)',
     '  1-on-1 meetings (w=0.15, x0=4)', '  Recognitions given (w=0.15, x0=3)',
     '  Team contributions (w=0.20, x0=2)'],
    'CIS = WeightedMean(6 sigmoid-normalized channels)\nExcluded if total collab < 12 activities',
    ['Output: 0-100', '80+: Highly collaborative', '50-79: Active collaborator', '< 50: Could collaborate more',
     'Excluded (redistributed) if too sparse'],
    'Weight: 10% of final CPIS score'
)

formula_slide(
    'D5', 'CPIS Dimension 5: Consistency & Reliability Index (CRI) — 10%',
    'Measures reliability: on-time delivery, steady velocity (no wild swings), streak of active days, consistent ratings, and deadline adherence.',
    ['OnTimeRate = on-time deliveries / total deliveries * 100', 'VelocityConsistency = max(0, (1 - variance/50) * 100)',
     'StreakFactor = sigmoid(streakDays, k=0.1, x0=14) * 100', 'RatingConsistency = max(0, (1 - stddev/2) * 100)',
     'DeadlineScore = (total - missed) / total * 100'],
    'CRI = 0.30*OnTime + 0.25*Velocity + 0.20*Streak + 0.15*Rating + 0.10*Deadline',
    ['Output: 0-100', '85+: Extremely reliable', '70-84: Dependable', '50-69: Somewhat consistent',
     '< 50: Inconsistent (flag for coaching)'],
    'Weight: 10% of final CPIS score'
)

formula_slide(
    'D6', 'CPIS Dimension 6: Growth Trajectory Score (GTS) — 8%',
    'Measures improvement over time. Uses linear regression on historical scores to detect if someone is getting better, staying flat, or declining. Also tracks skills, training, and development plans.',
    ['TrendScore = sigmoid(regression_slope, k=0.5, x0=0) * 100', 'SkillGrowth = sigmoid(progressions, k=0.5, x0=3) * 100',
     'TrainingScore = sigmoid(completions, k=0.5, x0=2) * 100', 'DevPlanProgress = direct 0-100',
     'ReadinessScore = promotion readiness 0-100'],
    'GTS = 0.35*Trend + 0.20*Skill + 0.15*Training + 0.15*DevPlan + 0.15*Readiness',
    ['Output: 0-100', 'slope > 0: Improving (score rises)', 'slope = 0: Stable', 'slope < 0: Declining (concerning)',
     'Default 50 if no historical data'],
    'Weight: 8% of final CPIS score'
)

formula_slide(
    'D7', 'CPIS Dimension 7: Evidence Quality Score (EQS) — 8%',
    'Scores the quality and diversity of work evidence submitted. Rewards verified evidence with measurable impact. Encourages submitting different types of proof.',
    ['VerificationRate = verified / total * 100', 'AvgImpact = mean of impact scores (0-100)',
     'AvgQuality = mean of quality scores (0-100)', 'DiversityBonus = sigmoid(types_count, k=0.5, x0=3) * 100'],
    'EQS = 0.25*Verification + 0.30*Impact + 0.25*Quality + 0.20*Diversity',
    ['Output: 0-100', 'Default: 50 if no evidence (neutral)', '80+: Strong, verified, diverse evidence', '50-79: Some evidence submitted',
     '< 50: Weak evidence portfolio'],
    'Weight: 8% of final CPIS score'
)

formula_slide(
    'D8', 'CPIS Dimension 8: Initiative & Innovation Index (III) — 7%',
    'Rewards going above and beyond: innovation contributions, mentoring, knowledge sharing, process improvements, and voluntary goals. Each measured via sigmoid curves.',
    ['Innovation = sigmoid(contributions, k=0.6, x0=2) * 100', 'Mentoring = sigmoid(sessions, k=0.4, x0=3) * 100',
     'Knowledge = sigmoid(sharing, k=0.5, x0=2) * 100', 'Process = sigmoid(improvements, k=0.6, x0=1) * 100',
     'Voluntary = sigmoid(goals, k=0.5, x0=2) * 100'],
    'III = 0.25*Innovation + 0.20*Mentoring + 0.20*Knowledge + 0.15*Process + 0.20*Voluntary',
    ['Output: 0-100', 'Excluded if no initiative data tracked', '80+: Highly innovative, proactive contributor', '50-79: Shows initiative',
     '< 50: Could take more ownership'],
    'Weight: 7% of final CPIS score'
)


# ═══════════════════════════════════════════════════════════
# FINAL CPIS ASSEMBLY
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             'CPIS Final Score Assembly', 28, VIOLET, True)

# Step 1-6
steps_data = [
    ('Step 1', 'Compute 8 Dimensions', 'Calculate GAI, RQS, FSI, CIS, CRI, GTS, EQS, III using formulas on previous slides', TEAL),
    ('Step 2', 'Dynamic Re-weighting', 'If a dimension has no data, exclude it and redistribute its weight proportionally to active dimensions', AMBER),
    ('Step 3', 'Tenure Factor', 'TenureFactor = min(1.12, 1 + years * 0.025) — up to 12% bonus at 4.8 years', GREEN),
    ('Step 4', 'Raw Score', 'RawScore = min(100, Sum(dim_i * weight_i)) * TenureFactor', WHITE),
    ('Step 5', 'Fairness Adjustment', 'Bayesian shrinkage for sparse data + disparate impact check + bias flag review', AMBER),
    ('Step 6', 'Final Score', 'FinalScore = clamp(RawScore + FairnessAdj, 0, 100)', GREEN),
]

for i, (step, title, desc, color) in enumerate(steps_data):
    y = Inches(1.2 + i * 0.75)
    add_rect(slide, Inches(0.5), y, Inches(6.0), Inches(0.65), SLATE)
    add_text_box(slide, Inches(0.6), y + Inches(0.02), Inches(0.8), Inches(0.3), step, 11, color, True)
    add_text_box(slide, Inches(1.5), y + Inches(0.02), Inches(2.0), Inches(0.3), title, 12, WHITE, True)
    add_text_box(slide, Inches(0.6), y + Inches(0.3), Inches(5.8), Inches(0.3), desc, 10, LIGHT)

# Right side: output table
add_rect(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.5), SLATE, VIOLET)
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5.4), Inches(0.4), 'FINAL OUTPUT', 16, VIOLET, True)

outputs = [
    ('Score', '0-100 (e.g., 78.5)'),
    ('Grade', 'A+=95+, A=85+, B+=78+, B=70+, C+=62+, C=50+, D=35+, F=<35'),
    ('Stars', '5=90+, 4=75+, 3=55+, 2=35+, 1=<35'),
    ('Rank', 'Exceptional(95+), Elite(85+), High Achiever(75+), Strong(65+), Solid(55+), Developing(45+), Emerging(35+), Needs Support(<35)'),
    ('Confidence', 'Margin = max(2, 15*(1-confidence)) — narrow when data-rich, wide when sparse'),
    ('Growth', 'Improving / Stable / Declining — from linear regression slope'),
]

for i, (label, desc) in enumerate(outputs):
    y = Inches(1.8 + i * 0.7)
    add_text_box(slide, Inches(7.3), y, Inches(1.5), Inches(0.3), label, 13, TEAL, True)
    add_text_box(slide, Inches(8.8), y, Inches(3.8), Inches(0.6), desc, 10.5, LIGHT)


# ═══════════════════════════════════════════════════════════
# CLOSING SLIDE
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1),
             'Every Number is Mathematically Derived', 32, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(0.6),
             'Zero Guesswork  |  Zero Hardcoded Values  |  Fully Auditable', 20, TEAL, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.6),
             '17 Core Formulas  +  6 Composite Systems  +  8 CPIS Dimensions', 18, AMBER, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.6),
             'Missing data = neutral (50), not penalizing. Bias detected = flagged. Legally defensible.', 14, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.4),
             'PMS Platform  |  Confidential  |  March 2026', 12, GRAY, False, PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), 'PMS_Backend_Formulas.pptx')
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
